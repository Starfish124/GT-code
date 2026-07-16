"""Document tools: Excel, PowerPoint and Word creation.

Kept forgiving on purpose — small local models produce slightly different JSON
shapes, so every tool accepts both the documented shape and common variations
(a bare list of rows, strings instead of dicts, …) rather than erroring out.

Dependencies (openpyxl / python-pptx / python-docx) are imported lazily so GT
still starts on a machine where they failed to install; the tool just returns
an actionable install hint instead.
"""

import sys

from .base import Tool

# Addressed to the MODEL, with the exact recovery step: a live session showed
# a 3B respond to the old bare hint by inventing a 'pip' tool and hallucinating
# successful installs, instead of using run_command with GT's own python.
_INSTALL_HINT = ('ERROR: the \'{pkg}\' package is not installed in GT\'s '
                 'environment, so this tool cannot run yet. To fix it, call '
                 'the run_command tool ONCE with exactly this command: '
                 '"{py}" -m pip install {pkg} — then, if it succeeded, retry '
                 'this tool ONCE. If the install fails, stop and report it '
                 '(the user can re-run setup.bat / setup.sh).')


def _install_hint(pkg):
    return _INSTALL_HINT.format(pkg=pkg, py=sys.executable)


def _rows_of(sheet: dict) -> list:
    rows = sheet.get("rows") or []
    headers = sheet.get("headers")
    return ([headers] + list(rows)) if headers else list(rows)


def _col_index(ident, headers):
    """Resolve a chart column identifier — a header NAME, a column letter, or a
    1-based index — to a 1-based column number, or None if it can't be placed."""
    if isinstance(ident, bool):          # bool is an int subclass; reject it
        return None
    if isinstance(ident, int):
        return ident if ident >= 1 else None
    s = str(ident or "").strip()
    if not s:
        return None
    low = [str(h).strip().lower() for h in (headers or [])]
    if s.lower() in low:                 # header-name match (what a model writes)
        return low.index(s.lower()) + 1
    if s.isdigit():
        return int(s) or None
    if s.isalpha() and len(s) <= 3:      # a column letter like "B"
        try:
            from openpyxl.utils import column_index_from_string
            return column_index_from_string(s.upper())
        except Exception:
            return None
    return None


def _add_chart(ws, spec, headers):
    """Add one native chart to a worksheet from an optional {"chart": ...} spec.

    Kept forgiving and non-fatal: a malformed spec is skipped, never allowed to
    lose the data workbook the model actually needed. Columns are named by their
    header text (categories = the label column, values = one or more numeric
    columns), which is what a small model can produce from what it just read."""
    from openpyxl.chart import BarChart, LineChart, PieChart, Reference
    if not isinstance(spec, dict):
        return
    kinds = {"bar": BarChart, "line": LineChart, "pie": PieChart}
    chart = kinds.get(str(spec.get("type", "bar")).lower(), BarChart)()
    if spec.get("title"):
        chart.title = str(spec["title"])
    has_header = bool(headers)
    first_data = 2 if has_header else 1
    if ws.max_row < first_data:
        return
    val_ids = spec.get("values")
    val_ids = val_ids if isinstance(val_ids, list) else [val_ids]
    val_cols = [c for c in (_col_index(v, headers) for v in val_ids) if c]
    if not val_cols:
        return
    data = Reference(ws, min_col=min(val_cols), max_col=max(val_cols),
                     min_row=1 if has_header else first_data, max_row=ws.max_row)
    chart.add_data(data, titles_from_data=has_header)
    cat_col = _col_index(spec.get("categories"), headers)
    if cat_col:
        chart.set_categories(Reference(ws, min_col=cat_col, max_col=cat_col,
                                       min_row=first_data, max_row=ws.max_row))
    from openpyxl.utils import get_column_letter
    anchor = str(spec.get("anchor")
                 or f"{get_column_letter(ws.max_column + 2)}2")
    ws.add_chart(chart, anchor)


class CreateExcel(Tool):
    name = "create_excel"
    description = ("Create an .xlsx Excel workbook. Headers are bolded and "
                   "column widths auto-fit. A sheet may include an optional "
                   "chart built from its own columns.")
    args = {
        "path": "Output file path ending in .xlsx.",
        "sheets": ('List of sheets: [{"name": "Sheet1", "headers": ["col", ...], '
                   '"rows": [[cell, ...], ...], "chart": {"type": "bar|line|pie", '
                   '"title": "...", "categories": "<header>", "values": '
                   '"<header>" or ["<header>", ...]}}]  (chart is optional)'),
    }
    arg_types = {"sheets": {"type": "array", "items": {"type": "object"}}}
    required = ("path", "sheets")
    changes_system = True

    def run(self, args, ctx):
        try:
            from openpyxl import Workbook
            from openpyxl.styles import Font
            from openpyxl.utils import get_column_letter
        except ImportError:
            return _install_hint("openpyxl")

        sheets = args.get("sheets") or []
        if isinstance(sheets, dict):
            sheets = [sheets]
        if not sheets and args.get("rows"):        # bare rows at top level
            sheets = [{"name": "Sheet1",
                       "headers": args.get("headers"), "rows": args["rows"]}]
        if not sheets:
            return "ERROR: 'sheets' is empty — pass at least one sheet with rows."

        p = ctx.resolve(args.get("path", "workbook.xlsx"))
        summary = ", ".join(f"{s.get('name', f'Sheet{i+1}')} "
                            f"({len(s.get('rows') or [])} rows)"
                            for i, s in enumerate(sheets))
        if not ctx.approve(f"Create Excel {p.name}", f"{p}\nsheets: {summary}",
                           key="docs"):
            return "DENIED: user declined creating the workbook."

        wb = Workbook()
        wb.remove(wb.active)
        for i, sheet in enumerate(sheets):
            if isinstance(sheet, list):            # model sent a bare row list
                sheet = {"rows": sheet}
            ws = wb.create_sheet(str(sheet.get("name") or f"Sheet{i+1}")[:31])
            has_header = bool(sheet.get("headers"))
            for row in _rows_of(sheet):
                ws.append(list(row) if isinstance(row, (list, tuple)) else [row])
            if has_header and ws.max_row >= 1:
                for cell in ws[1]:
                    cell.font = Font(bold=True)
                ws.freeze_panes = "A2"
            for col in range(1, ws.max_column + 1):
                width = max((len(str(ws.cell(r, col).value or ""))
                             for r in range(1, ws.max_row + 1)), default=8)
                ws.column_dimensions[get_column_letter(col)].width = \
                    min(max(width + 2, 8), 60)
            if isinstance(sheet, dict) and sheet.get("chart"):
                try:
                    _add_chart(ws, sheet["chart"], sheet.get("headers"))
                except Exception:
                    pass          # a chart glitch must never lose the workbook
        try:
            p.parent.mkdir(parents=True, exist_ok=True)
            wb.save(p)
        except Exception as e:
            return f"ERROR saving {p}: {e}"
        nchart = sum(1 for s in sheets
                     if isinstance(s, dict) and s.get("chart"))
        tail = f", {nchart} chart(s)" if nchart else ""
        return f"OK: created {p} with {len(sheets)} sheet(s){tail}."


class CreatePowerPoint(Tool):
    name = "create_powerpoint"
    description = ("Create a .pptx PowerPoint deck: an optional title slide "
                   "plus title+bullet slides (with optional speaker notes).")
    args = {
        "path": "Output file path ending in .pptx.",
        "title": "Deck title for the opening slide (optional).",
        "subtitle": "Subtitle for the opening slide (optional).",
        "slides": ('List of slides: [{"title": "...", "bullets": ["...", ...], '
                   '"notes": "optional speaker notes"}, ...]'),
    }
    arg_types = {"slides": {"type": "array", "items": {"type": "object"}}}
    required = ("path", "slides")
    changes_system = True

    def run(self, args, ctx):
        try:
            from pptx import Presentation
        except ImportError:
            return _install_hint("python-pptx")

        slides = args.get("slides") or []
        if isinstance(slides, dict):
            slides = [slides]
        if not slides and not args.get("title"):
            return "ERROR: pass 'slides' (and optionally a deck 'title')."

        p = ctx.resolve(args.get("path", "deck.pptx"))
        outline = "\n".join(f"  {i+1}. {s.get('title', '(untitled)')}"
                            if isinstance(s, dict) else f"  {i+1}. {s}"
                            for i, s in enumerate(slides))
        if not ctx.approve(f"Create PowerPoint {p.name}",
                           f"{p}\n{len(slides)} slide(s):\n{outline}",
                           key="docs"):
            return "DENIED: user declined creating the deck."

        prs = Presentation()
        if args.get("title"):
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = str(args["title"])
            if args.get("subtitle") and len(slide.placeholders) > 1:
                slide.placeholders[1].text = str(args["subtitle"])

        for s in slides:
            if isinstance(s, str):
                s = {"title": s}
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = str(s.get("title", ""))
            bullets = s.get("bullets") or []
            if bullets:
                tf = slide.placeholders[1].text_frame
                tf.text = str(bullets[0])
                for b in bullets[1:]:
                    tf.add_paragraph().text = str(b)
            if s.get("notes"):
                slide.notes_slide.notes_text_frame.text = str(s["notes"])
        try:
            p.parent.mkdir(parents=True, exist_ok=True)
            prs.save(p)
        except Exception as e:
            return f"ERROR saving {p}: {e}"
        return f"OK: created {p} with {len(prs.slides)} slide(s)."


class CreateWord(Tool):
    name = "create_word"
    description = ("Create a .docx Word document from a list of blocks "
                   "(headings, paragraphs, bullet lists).")
    args = {
        "path": "Output file path ending in .docx.",
        "blocks": ('List of blocks: [{"type": "heading|paragraph|bullets", '
                   '"text": "...", "items": ["..."], "level": 1}, ...] — '
                   'plain strings are treated as paragraphs.'),
    }
    arg_types = {"blocks": {"type": "array"}}
    required = ("path", "blocks")
    changes_system = True

    def run(self, args, ctx):
        try:
            from docx import Document
        except ImportError:
            return _install_hint("python-docx")

        blocks = args.get("blocks") or []
        if isinstance(blocks, (str, dict)):
            blocks = [blocks]
        if not blocks:
            return "ERROR: 'blocks' is empty."

        p = ctx.resolve(args.get("path", "document.docx"))
        if not ctx.approve(f"Create Word doc {p.name}",
                           f"{p}\n{len(blocks)} block(s)", key="docs"):
            return "DENIED: user declined creating the document."

        doc = Document()
        for b in blocks:
            if isinstance(b, str):
                doc.add_paragraph(b)
                continue
            kind = (b.get("type") or "paragraph").lower()
            if kind == "heading":
                doc.add_heading(str(b.get("text", "")),
                                level=min(max(int(b.get("level", 1)), 0), 9))
            elif kind == "bullets":
                for item in (b.get("items") or []):
                    doc.add_paragraph(str(item), style="List Bullet")
            else:
                doc.add_paragraph(str(b.get("text", "")))
        try:
            p.parent.mkdir(parents=True, exist_ok=True)
            doc.save(p)
        except Exception as e:
            return f"ERROR saving {p}: {e}"
        return f"OK: created {p} with {len(blocks)} block(s)."


OFFICE_TOOLS = [CreateExcel(), CreatePowerPoint(), CreateWord()]
